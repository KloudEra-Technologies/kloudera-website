import os
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 1. Colors Setup (Brand Identity)
C_DARK_BLUE = (10, 44, 109)    # #0A2C6D (Primary Dark Blue)
C_WHITE = (255, 255, 255)      # White
C_LIGHT_BLUE = (59, 130, 246)  # #3B82F6 (Light Blue)
C_ELEC_BLUE = (37, 99, 235)   # #2563EB (Electric Blue)
C_CYAN = (6, 182, 212)         # #06B6D4 (Accent Cyan)
C_TEXT_DARK = (15, 23, 42)     # #0F172A (Text Dark Slate)
C_TEXT_MUTED = (100, 116, 139) # #64748B (Text Muted Slate)

# Font name
FONT_TITLE = 'Arial'
FONT_BODY = 'Arial'

# Logo Path
LOGO_PATH = 'KlouderaNewLogo-removebg-preview.png'
LOGO_WHITE_PATH = 'KlouderaNewLogo-white.png'
CEO1_PATH = 'CEO (1).jpeg'
CEO2_PATH = 'CEO (2).jpeg'

# Load website content json
def load_content():
    content_path = os.path.join('src', 'data', 'website_content.json')
    if os.path.exists(content_path):
        with open(content_path, 'r', encoding='utf-8') as f:
            return json.load(f)
    return {}

# Slide Layout Helper - Dark (Now Light/White to make Logo fully visible)
def add_dark_slide(prs, title_text, subtitle_text=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255) # Pure White
    
    # Accent cyan bar at the top (to keep it premium and consistent)
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.12))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = RGBColor(*C_CYAN)
    top_bar.line.fill.background()
    
    # Add logo if exists (original logo is used directly)
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, Inches(10.2), Inches(0.4), width=Inches(2.5))
        
    # Title
    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(2.8), Inches(11.333), Inches(1.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*C_DARK_BLUE)
    p.font.name = FONT_TITLE
    
    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(*C_ELEC_BLUE)
        p2.font.name = FONT_TITLE
        p2.space_before = Pt(20)
        
    return slide

# Slide Layout Helper - Content (White)
def add_content_slide(prs, title_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*C_WHITE)
    
    # Accent cyan bar at the top
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.12))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = RGBColor(*C_CYAN)
    top_bar.line.fill.background()
    
    # Add logo if exists
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, Inches(10.2), Inches(0.4), width=Inches(2.5))
        
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(9.0), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Inches(0.0)
    p = tf.paragraphs[0]
    p.text = title_text.upper()
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*C_DARK_BLUE)
    p.font.name = FONT_TITLE
    
    return slide

# TextBox Helper
def add_textbox(slide, left, top, width, height, text, size=13, bold=False, color=C_TEXT_DARK, align=PP_ALIGN.LEFT, italic=False):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Inches(0.0)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = RGBColor(*color)
    p.font.name = FONT_BODY
    p.alignment = align
    return tx, tf

# Card Helper
def add_card(slide, left, top, width, height, bg_color=C_WHITE, border_color=(226, 232, 240)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*bg_color)
    if border_color:
        shape.line.color.rgb = RGBColor(*border_color)
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def generate_deck():
    content = load_content()
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # -------------------------------------------------------------------------
    # SLIDE 1: Cover
    # -------------------------------------------------------------------------
    title = "CII STARTUPRENEUR AWARDS 2026"
    subtitle = "Kloudera Technologies Pvt. Ltd.\nEmpowering Secure Digital Transformation"
    slide = add_dark_slide(prs, title, subtitle)
    
    # Add a glowing element or futuristic design (subtle shape)
    glow = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1), Inches(4.5), Inches(5), Inches(4))
    glow.fill.solid()
    glow.fill.fore_color.rgb = RGBColor(239, 246, 255) # Light Blue 50
    glow.line.fill.background()
    
    # Info note
    add_textbox(slide, Inches(1.0), Inches(6.0), Inches(11.333), Inches(1.0),
                "Official Submission | Category: Technology & Digital Innovation\nwww.kloudera.ai",
                size=12, color=C_TEXT_MUTED)
    
    slide.notes_slide.notes_text_frame.text = (
        "Welcome judges to the presentation of Kloudera Technologies for the CII Startupreneur Awards 2026. "
        "Kloudera Technologies is a high-assurance enterprise platform architecting secure, Zero-Trust digital enterprises. "
        "This presentation will showcase our innovations, technology stack, social impact, and strategic roadmap."
    )

    # -------------------------------------------------------------------------
    # SLIDE 2: Company Overview
    # -------------------------------------------------------------------------
    slide = add_content_slide(prs, "Company Overview")
    
    # Left Block: Vision/Mission (M365-like design card)
    add_card(slide, Inches(0.8), Inches(1.5), Inches(5.6), Inches(4.8), bg_color=(248, 250, 252))
    add_textbox(slide, Inches(1.2), Inches(1.8), Inches(4.8), Inches(0.4), "VISION", size=14, bold=True, color=C_ELEC_BLUE)
    add_textbox(slide, Inches(1.2), Inches(2.2), Inches(4.8), Inches(1.0), 
                "Built for High-Stakes Enterprise Resilience.", size=18, bold=True, color=C_TEXT_DARK)
    
    add_textbox(slide, Inches(1.2), Inches(3.4), Inches(4.8), Inches(0.4), "MISSION", size=14, bold=True, color=C_ELEC_BLUE)
    add_textbox(slide, Inches(1.2), Inches(3.8), Inches(4.8), Inches(1.2), 
                content.get("about", {}).get("missionDesc", "We believe in your success and that technology can drive the best results for your business, no matter your industry or goals."), 
                size=14, color=C_TEXT_MUTED)

    # Right Block: Core Capabilities & Timeline/Details
    add_card(slide, Inches(6.8), Inches(1.5), Inches(5.6), Inches(4.8), bg_color=C_WHITE)
    add_textbox(slide, Inches(7.2), Inches(1.8), Inches(4.8), Inches(0.4), "GOVERNMENT ACCREDITATION", size=14, bold=True, color=C_CYAN)
    
    # MSME & Startup India
    add_textbox(slide, Inches(7.2), Inches(2.3), Inches(4.8), Inches(0.8), 
                "• #startupindia Recognized\n  DPIIT Registered Venture (Govt. of India)", size=13, bold=True, color=C_TEXT_DARK)
    add_textbox(slide, Inches(7.2), Inches(3.1), Inches(4.8), Inches(0.8), 
                "• MSME Certified Enterprise\n  Ministry of Micro, Small & Medium Enterprises", size=13, bold=True, color=C_TEXT_DARK)
    
    add_textbox(slide, Inches(7.2), Inches(4.0), Inches(4.8), Inches(0.4), "OFFICE LOCATIONS", size=14, bold=True, color=C_CYAN)
    add_textbox(slide, Inches(7.2), Inches(4.5), Inches(4.8), Inches(1.2), 
                "• Bengaluru R&D HQ (Karnataka)\n• Pune Operations Center (Maharashtra)", size=13, color=C_TEXT_DARK)
                
    slide.notes_slide.notes_text_frame.text = (
        "Kloudera Technologies is registered under the Ministry of MSME and recognized under the flagship Startup India initiative. "
        "Our R&D headquarters is in Bengaluru and our main operations center is in Pune. "
        "Our vision is to provide high-stakes enterprise resilience for global multinationals and Indian businesses alike."
    )

    # -------------------------------------------------------------------------
    # SLIDE 3: Problem Statement
    # -------------------------------------------------------------------------
    slide = add_content_slide(prs, "The Enterprise Challenge")
    
    # 3x2 Grid for 6 major challenges
    challenges = [
        ("CYBER SECURITY THREATS", "Ransomware and sophisticated phishing campaigns target data and infrastructure daily."),
        ("LEGACY INFRASTRUCTURE", "Rigid on-premise components limit performance, scaling speed, and business agility."),
        ("AI ADOPTION HURDLES", "Difficulty in designing pipelines, securing ML models (MLSecOps), and preventing data leakage."),
        ("CLOUD GOVERNANCE GAP", "Inefficient multi-cloud operations lead to major cost waste and security misconfigurations."),
        ("COMPLIANCE BURDEN", "Strict regulatory standards (GRC, GDPR, DPDP) require constant posture audits."),
        ("SPECIALIZED SKILL DRAIN", "Acute shortage of top-tier talent in advanced cloud architectures, security, and AI.")
    ]
    
    col_w = Inches(3.6)
    row_h = Inches(2.0)
    for i, (title_ch, desc_ch) in enumerate(challenges):
        col = i % 3
        row = i // 3
        left = Inches(0.8) + col * (col_w + Inches(0.4))
        top = Inches(1.8) + row * (row_h + Inches(0.4))
        
        add_card(slide, left, top, col_w, row_h, bg_color=(254, 254, 254))
        # Small icon box indicator (colored line on left side)
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.08), row_h)
        accent.fill.solid()
        accent.fill.fore_color.rgb = RGBColor(*C_CYAN if i%2==0 else C_LIGHT_BLUE)
        accent.line.fill.background()
        
        add_textbox(slide, left + Inches(0.2), top + Inches(0.2), col_w - Inches(0.4), Inches(0.4),
                    title_ch, size=11, bold=True, color=C_TEXT_DARK)
        add_textbox(slide, left + Inches(0.2), top + Inches(0.7), col_w - Inches(0.4), Inches(1.1),
                    desc_ch, size=11, color=C_TEXT_MUTED)
                    
    slide.notes_slide.notes_text_frame.text = (
        "Enterprise IT leaders face a combined threat: increasing cyber attacks, legacy systems, complex compliance, "
        "and a major skill gap. Kloudera addresses these interconnected issues through a unified platform approach."
    )

    # -------------------------------------------------------------------------
    # SLIDE 4: Our Solution
    # -------------------------------------------------------------------------
    slide = add_content_slide(prs, "Our Solution: Unified Resilience Ecosystem")
    
    # Text introduction on left
    add_textbox(slide, Inches(0.8), Inches(1.8), Inches(4.5), Inches(4.5),
                "THE SOLUTIONS LAYER\n\n"
                "Kloudera operates as a unified platform bridging cybersecurity, cloud operations, automation, AI enablement, and high-performance hardware.\n\n"
                "Rather than hiring disconnected vendors, enterprise clients receive a consolidated, secure digital core built on Zero-Trust architecture.",
                size=14, color=C_TEXT_DARK)
    
    # Graphic on right: 5 rounded cards vertically stacked or clustered with connections
    # Let's do horizontal/vertical list with nice offsets representing layers
    layers = [
        ("01 / ARTIFICIAL INTELLIGENCE", "Gen-AI pipelines, custom LLMs & MLSecOps frameworks."),
        ("02 / CLOUD COMPUTING", "Multi-cloud migration, orchestration & compliance-ready serverless."),
        ("03 / CYBERSECURITY OPERATIONS", "vCISO governance, 24/7 Security surveillance (SoC), VAPT."),
        ("04 / INTELLIGENT AUTOMATION", "Robotic Process Automation (RPA) and DevOps CI/CD pipelines."),
        ("05 / DIGITAL ENGINEERING", "Custom microservices app development and database optimization.")
    ]
    
    for i, (ly_title, ly_desc) in enumerate(layers):
        left_pos = Inches(6.0) + Inches(i * 0.25)
        top_pos = Inches(1.5) + Inches(i * 1.05)
        card_w = Inches(6.0)
        card_h = Inches(0.9)
        
        # Glow color based on index
        color_grad = C_CYAN if i==0 else (C_LIGHT_BLUE if i==1 else (C_ELEC_BLUE if i==2 else C_DARK_BLUE))
        
        add_card(slide, left_pos, top_pos, card_w, card_h, bg_color=(250, 252, 255))
        # Left bar indicator
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_pos, top_pos, Inches(0.12), card_h)
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(*color_grad)
        bar.line.fill.background()
        
        add_textbox(slide, left_pos + Inches(0.3), top_pos + Inches(0.12), card_w - Inches(0.5), Inches(0.3),
                    ly_title, size=11, bold=True, color=C_TEXT_DARK)
        add_textbox(slide, left_pos + Inches(0.3), top_pos + Inches(0.42), card_w - Inches(0.5), Inches(0.45),
                    ly_desc, size=11, color=C_TEXT_MUTED)
                    
    slide.notes_slide.notes_text_frame.text = (
        "We unite five critical technology pillars. By layering cybersecurity on top of cloud, AI, and automation, "
        "we eliminate security gaps and dramatically speed up implementation times."
    )

    # -------------------------------------------------------------------------
    # SLIDE 5: Service Portfolio
    # -------------------------------------------------------------------------
    slide = add_content_slide(prs, "Comprehensive Service Portfolio")
    
    # 4 columns of cards
    services = [
        ("CYBERSECURITY", "vCISO & vDPO\nVAPT / Red Team\nEDR, XDR, MDR\n24/7 SoC Service\nCloud GRC & Security\nMLSecOps (AI Security)"),
        ("CLOUD & INFRA", "Cloud Migrations\nMulti-Cloud Setup\nServerless Compute\nCost Optimization\nInfra Management\nDR as a Service"),
        ("AI & AUTOMATION", "Gen-AI Pipelines\nCustom LLM Rigs\nRPA Automation\nDevOps & CI/CD\nIT Automation\nInfrastructure as Code"),
        ("CONSULTING & DEV", "Digital Consulting\nCustom Software Dev\nMicroservices Arch\nDatabase Design\nSAP, Oracle, CRM\nStaff Augmentation")
    ]
    
    col_w = Inches(2.7)
    for i, (title_sv, items_sv) in enumerate(services):
        left = Inches(0.8) + i * (col_w + Inches(0.3))
        top = Inches(1.6)
        card_h = Inches(4.8)
        
        add_card(slide, left, top, col_w, card_h, bg_color=C_WHITE)
        # Header block of card
        hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, col_w, Inches(0.8))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = RGBColor(*C_DARK_BLUE)
        hdr.line.fill.background()
        
        add_textbox(slide, left + Inches(0.2), top + Inches(0.25), col_w - Inches(0.4), Inches(0.5),
                    title_sv, size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        
        add_textbox(slide, left + Inches(0.2), top + Inches(1.1), col_w - Inches(0.4), Inches(3.4),
                    items_sv, size=11, color=C_TEXT_DARK, align=PP_ALIGN.LEFT)
                    
    slide.notes_slide.notes_text_frame.text = (
        "This slide presents a granular list of our capabilities. "
        "It includes managed cybersecurity services, cloud governance, automation pipelines, and custom enterprise software development."
    )

    # -------------------------------------------------------------------------
    # SLIDE 6: Technology Stack
    # -------------------------------------------------------------------------
    slide = add_content_slide(prs, "Enterprise Technology Stack")
    
    # 4 columns for stack categories
    tech_stacks = [
        ("CLOUD PROVIDERS", "Microsoft Azure\nAmazon Web Services\nGoogle Cloud Platform\nHybrid Datacenters"),
        ("SECURITY & COMPLIANCE", "Microsoft Entra\nIAM / IGA / PAM\nEDR / XDR / VMDR\nSprinto (SOC 2, ISO)\nCross Cipher Threat Intelligence"),
        ("AI & AUTOMATION", "Large Language Models (LLMs)\nTensorFlow / PyTorch\nUiPath (RPA)\nTerraform (IaC)\nDevOps Pipelines"),
        ("DIGITAL PLATFORMS", "Node.js / Next.js\nMicroservices (Docker, K8s)\nMongoDB / SQL\nSAP / Oracle ERP\nSalesforce CRM")
    ]
    
    col_w = Inches(2.7)
    for i, (stack_name, stack_details) in enumerate(tech_stacks):
        left = Inches(0.8) + i * (col_w + Inches(0.3))
        top = Inches(1.8)
        card_h = Inches(4.5)
        
        add_card(slide, left, top, col_w, card_h, bg_color=(245, 248, 253))
        # Small accent block at the top of the card
        acc = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, col_w, Inches(0.1))
        acc.fill.solid()
        acc.fill.fore_color.rgb = RGBColor(*C_CYAN if i%2==0 else C_LIGHT_BLUE)
        acc.line.fill.background()
        
        add_textbox(slide, left + Inches(0.2), top + Inches(0.3), col_w - Inches(0.4), Inches(0.6),
                    stack_name, size=11, bold=True, color=C_DARK_BLUE)
        add_textbox(slide, left + Inches(0.2), top + Inches(1.1), col_w - Inches(0.4), Inches(3.2),
                    stack_details, size=11, color=C_TEXT_DARK)
                    
    slide.notes_slide.notes_text_frame.text = (
        "Our technology stack is built on industry-standard global technologies. We leverage Microsoft Azure, AWS, "
        "and Google Cloud for cloud hosting, and use advanced tools like Terraform, Sprinto, and proprietary MLSecOps to ensure compliance."
    )

    # -------------------------------------------------------------------------
    # SLIDE 7: Innovation
    # -------------------------------------------------------------------------
    slide = add_content_slide(prs, "Key Innovations & Market Edge")
    
    # 3 major innovations in visual columns
    innovations = [
        ("AI-FIRST ARCHITECTURE", "Center of Excellence", 
         "We build bespoke LLM rigs and pipeline engines tailored to enterprise workflows, bypassing generic, high-cost public API endpoints."),
        ("SECURE-BY-DESIGN CORE", "Zero-Trust Engineering", 
         "Cybersecurity isn't an afterthought. Every cloud migration, custom app, and automated script has built-in IAM, GRC, and encryption."),
        ("VENDOR AGNOSTIC CLOUD", "Hybrid Orchestration", 
         "We deploy multi-cloud systems integrating Microsoft, AWS, MongoDB, and partner security solutions for optimal resilience.")
    ]
    
    col_w = Inches(3.6)
    for i, (inn_title, inn_sub, inn_desc) in enumerate(innovations):
        left = Inches(0.8) + i * (col_w + Inches(0.4))
        top = Inches(1.8)
        card_h = Inches(4.5)
        
        add_card(slide, left, top, col_w, card_h, bg_color=C_WHITE)
        # Gradient banner
        ban = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, col_w, Inches(1.0))
        ban.fill.solid()
        ban.fill.fore_color.rgb = RGBColor(*C_DARK_BLUE)
        ban.line.fill.background()
        
        add_textbox(slide, left + Inches(0.2), top + Inches(0.18), col_w - Inches(0.4), Inches(0.3),
                    inn_title, size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, left + Inches(0.2), top + Inches(0.55), col_w - Inches(0.4), Inches(0.3),
                    inn_sub, size=10, color=C_CYAN, align=PP_ALIGN.CENTER, italic=True)
        
        add_textbox(slide, left + Inches(0.25), top + Inches(1.3), col_w - Inches(0.5), Inches(2.9),
                    inn_desc, size=12, color=C_TEXT_DARK)
                    
    slide.notes_slide.notes_text_frame.text = (
        "Kloudera differentiates itself through an AI-First, secure-by-design, vendor-agnostic architecture. "
        "We allow clients to deploy customized LLMs locally, ensuring privacy and data ownership while optimizing cloud costs."
    )

    # -------------------------------------------------------------------------
    # SLIDE 8: Industries Served
    # -------------------------------------------------------------------------
    slide = add_content_slide(prs, "Industries Served")
    
    # 6 industries, 2 rows of 3 columns
    industries = [
        ("BFSI", "High-security financial structures, trading desks, capital management, and compliance checks (e.g. IndiaCapital, Fin Chikitsak)."),
        ("HEALTHCARE & LIFE SCIENCES", "Data privacy protection (vDPO), HIPAA compliance, secure medical records management."),
        ("RETAIL & E-COMMERCE", "High-volume secure transactions, microservices architecture, and cloud native customer experience scaling."),
        ("MANUFACTURING & OT", "Operational technology (OT) security, IoT sensor integration, and industrial process automation."),
        ("GOVERNMENT & PUBLIC", "Secure MSME digital integrations and national startup development initiatives."),
        ("EDUCATION & ACADEMICS", "Partnering with premier institutes like SIT Pune for technical research, talent pipelines, and labs.")
    ]
    
    col_w = Inches(3.6)
    row_h = Inches(2.0)
    for i, (ind_title, ind_desc) in enumerate(industries):
        col = i % 3
        row = i // 3
        left = Inches(0.8) + col * (col_w + Inches(0.4))
        top = Inches(1.8) + row * (row_h + Inches(0.4))
        
        add_card(slide, left, top, col_w, row_h, bg_color=(250, 252, 255))
        # Left line accent
        acc = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.08), row_h)
        acc.fill.solid()
        acc.fill.fore_color.rgb = RGBColor(*C_ELEC_BLUE if i%2==0 else C_CYAN)
        acc.line.fill.background()
        
        add_textbox(slide, left + Inches(0.2), top + Inches(0.2), col_w - Inches(0.4), Inches(0.4),
                    ind_title, size=12, bold=True, color=C_DARK_BLUE)
        add_textbox(slide, left + Inches(0.2), top + Inches(0.7), col_w - Inches(0.4), Inches(1.2),
                    ind_desc, size=11, color=C_TEXT_MUTED)
                    
    slide.notes_slide.notes_text_frame.text = (
        "We serve critical sectors like Banking, Healthcare, Manufacturing, and Education. "
        "Our clients include leading fintech companies like IndiaCapital and top universities like Symbiosis Institute of Technology."
    )

    # -------------------------------------------------------------------------
    # SLIDE 9: Market Opportunity
    # -------------------------------------------------------------------------
    slide = add_content_slide(prs, "Market Opportunity & Growth Vectors")
    
    # 4 market statistics boxes
    stats = [
        ("GEN-AI ENTERPRISE MARKET", "$1.3 Trillion", "Expected size by 2032 as companies build bespoke LLM models.", "Source: Bloomberg Intelligence"),
        ("CLOUD SERVICES SPENDING", "$679 Billion", "Global spending forecast by Gartner, driven by multi-cloud setups.", "Source: Gartner Research"),
        ("CYBERSECURITY & OT", "$266 Billion", "Rising threat landscapes demand proactive 24/7 security monitoring.", "Source: McKinsey & Co."),
        ("DIGITAL TRANSFORMATION", "$3.4 Trillion", "Global investment in cloud and automation workflows.", "Source: IDC Worldwide Spending Guide")
    ]
    
    col_w = Inches(2.7)
    for i, (st_name, st_val, st_desc, st_src) in enumerate(stats):
        left = Inches(0.8) + i * (col_w + Inches(0.3))
        top = Inches(1.8)
        card_h = Inches(4.5)
        
        add_card(slide, left, top, col_w, card_h, bg_color=C_WHITE)
        # Subtle gradient block for the number
        num_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, col_w, Inches(1.8))
        num_box.fill.solid()
        num_box.fill.fore_color.rgb = RGBColor(*C_DARK_BLUE)
        num_box.line.fill.background()
        
        add_textbox(slide, left + Inches(0.1), top + Inches(0.5), col_w - Inches(0.2), Inches(0.8),
                    st_val, size=24, bold=True, color=C_CYAN, align=PP_ALIGN.CENTER)
        
        add_textbox(slide, left + Inches(0.2), top + Inches(2.1), col_w - Inches(0.4), Inches(0.4),
                    st_name, size=11, bold=True, color=C_TEXT_DARK, align=PP_ALIGN.CENTER)
        add_textbox(slide, left + Inches(0.2), top + Inches(2.6), col_w - Inches(0.4), Inches(1.4),
                    st_desc, size=11, color=C_TEXT_MUTED, align=PP_ALIGN.CENTER)
        add_textbox(slide, left + Inches(0.2), top + Inches(4.1), col_w - Inches(0.4), Inches(0.3),
                    st_src, size=9, color=C_LIGHT_BLUE, align=PP_ALIGN.CENTER, italic=True)
                    
    slide.notes_slide.notes_text_frame.text = (
        "The market potential across Gen-AI, Cloud, Cybersecurity, and automation is massive. "
        "Bespoke LLMs alone are driving a $1.3 trillion transition, which aligns perfectly with our AI Center of Excellence."
    )

    # -------------------------------------------------------------------------
    # SLIDE 10: Business Model
    # -------------------------------------------------------------------------
    slide = add_content_slide(prs, "Scalable Business Model")
    
    # 4 vertical process cards
    phases = [
        ("01 / AUDIT & CONSULT", "Initial Assessment", "Perform risk-free infrastructure assessments, vCISO compliance audits, and security vulnerability testing (VAPT)."),
        ("02 / ARCHITECT & BUILD", "Custom Deployment", "Develop bespoke custom software, serverless cloud setups, Gen-AI pipelines, and high-performance hardware installations."),
        ("03 / MANAGED SERVICES", "Recurring Operations", "Deliver 24/7 Security Operations Center (SoC) surveillance, ongoing multi-cloud cost optimization, and proactive tech support."),
        ("04 / EXPAND & SECURE", "Continuous Success", "Provide software upgrades, regular compliance reporting, lifecycle hardware maintenance, and global expansion advisory.")
    ]
    
    col_w = Inches(2.7)
    for i, (ph_name, ph_sub, ph_desc) in enumerate(phases):
        left = Inches(0.8) + i * (col_w + Inches(0.3))
        top = Inches(1.8)
        card_h = Inches(4.5)
        
        add_card(slide, left, top, col_w, card_h, bg_color=(250, 252, 255))
        # Arrow header indicator
        hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, col_w, Inches(0.8))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = RGBColor(*C_DARK_BLUE if i%2==0 else C_ELEC_BLUE)
        hdr.line.fill.background()
        
        add_textbox(slide, left + Inches(0.1), top + Inches(0.25), col_w - Inches(0.2), Inches(0.4),
                    ph_name, size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        
        add_textbox(slide, left + Inches(0.2), top + Inches(1.1), col_w - Inches(0.4), Inches(0.4),
                    ph_sub, size=12, bold=True, color=C_CYAN, align=PP_ALIGN.CENTER)
        add_textbox(slide, left + Inches(0.25), top + Inches(1.6), col_w - Inches(0.5), Inches(2.6),
                    ph_desc, size=11, color=C_TEXT_DARK)
                    
    slide.notes_slide.notes_text_frame.text = (
        "Our business model drives long-term customer success. "
        "We acquire clients via initial audits, build their custom systems, and retain them through high-margin, recurring managed services."
    )

    # -------------------------------------------------------------------------
    # SLIDE 11: Customer Journey
    # -------------------------------------------------------------------------
    slide = add_content_slide(prs, "Customer Integration Journey")
    
    # Timeline chain (horizontal)
    steps = [
        ("1. DISCOVER", "Scope alignment & infrastructure mapping."),
        ("2. AUDIT", "Thorough VAPT & compliance threat detection."),
        ("3. DESIGN", "Architecting custom LLM rigs & cloud structures."),
        ("4. DEPLOY", "Secure migration & DevSecOps pipeline setup."),
        ("5. MONITOR", "Activate 24/7 Security SoC operations."),
        ("6. OPTIMIZE", "Ongoing cost control & tech scaling.")
    ]
    
    col_w = Inches(1.8)
    for i, (st_title, st_desc) in enumerate(steps):
        left = Inches(0.8) + i * (col_w + Inches(0.15))
        top = Inches(2.5)
        card_h = Inches(3.2)
        
        add_card(slide, left, top, col_w, card_h, bg_color=C_WHITE)
        # Circle step indicator on top
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.6), top - Inches(0.4), Inches(0.6), Inches(0.6))
        circ.fill.solid()
        circ.fill.fore_color.rgb = RGBColor(*C_CYAN)
        circ.line.fill.background()
        
        add_textbox(slide, left + Inches(0.6), top - Inches(0.3), Inches(0.6), Inches(0.4),
                    str(i+1), size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        
        add_textbox(slide, left + Inches(0.1), top + Inches(0.5), col_w - Inches(0.2), Inches(0.4),
                    st_title, size=11, bold=True, color=C_DARK_BLUE, align=PP_ALIGN.CENTER)
        add_textbox(slide, left + Inches(0.1), top + Inches(1.0), col_w - Inches(0.2), Inches(2.0),
                    st_desc, size=10, color=C_TEXT_MUTED, align=PP_ALIGN.CENTER)
                    
    # Horizontal connecting line behind circles
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.0), Inches(11.0), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(*C_CYAN)
    line.line.fill.background()
    
    slide.notes_slide.notes_text_frame.text = (
        "Trainees and enterprise clients follow a standard path. "
        "We move them from initial audit to production monitoring, ensuring they are protected and optimized at every step."
    )

    # -------------------------------------------------------------------------
    # SLIDE 12: Why Kloudera
    # -------------------------------------------------------------------------
    slide = add_content_slide(prs, "Why Kloudera Technologies?")
    
    # Comparison table (drawn manually via cards/grids for pixel-perfect design)
    # Header Row
    top_header = Inches(1.8)
    row_h = Inches(0.8)
    
    add_card(slide, Inches(0.8), top_header, Inches(2.5), row_h, bg_color=C_DARK_BLUE)
    add_textbox(slide, Inches(0.9), top_header + Inches(0.2), Inches(2.3), Inches(0.4), "CRITICAL CAPABILITY", size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    
    add_card(slide, Inches(3.4), top_header, Inches(4.5), row_h, bg_color=(241, 245, 249))
    add_textbox(slide, Inches(3.5), top_header + Inches(0.2), Inches(4.3), Inches(0.4), "TRADITIONAL IT VENDOR", size=12, bold=True, color=C_TEXT_DARK, align=PP_ALIGN.CENTER)
    
    add_card(slide, Inches(8.0), top_header, Inches(4.5), row_h, bg_color=C_ELEC_BLUE)
    add_textbox(slide, Inches(8.1), top_header + Inches(0.2), Inches(4.3), Inches(0.4), "KLOUDERA TECHNOLOGIES", size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    
    # Rows
    rows_data = [
        ("Security Standard", "Reactive, add-on software; frequent gaps.", "Zero-Trust default; native IAM & GRC monitoring."),
        ("AI Integration", "API wrappers, high security leakage risk.", "Bespoke local LLM rigs & Gen-AI pipelines."),
        ("Operational Speed", "Weeks of manual provisioning.", "Automation-first with Terraform IaC & CI/CD."),
        ("Cost Management", "Fixed rates, hidden cloud costs.", "Proactive multi-cloud optimization & scheduling."),
        ("Infrastructure Support", "Disconnected helpdesks.", "24/7 Operations Center (SoC) surveillance.")
    ]
    
    for idx, (cap, trad, kdr) in enumerate(rows_data):
        curr_top = top_header + Inches(0.9) + idx * Inches(0.85)
        
        # Capability
        add_card(slide, Inches(0.8), curr_top, Inches(2.5), Inches(0.8), bg_color=(248, 250, 252))
        add_textbox(slide, Inches(0.9), curr_top + Inches(0.2), Inches(2.3), Inches(0.4), cap, size=11, bold=True, color=C_DARK_BLUE, align=PP_ALIGN.CENTER)
        
        # Traditional
        add_card(slide, Inches(3.4), curr_top, Inches(4.5), Inches(0.8), bg_color=C_WHITE)
        add_textbox(slide, Inches(3.5), curr_top + Inches(0.2), Inches(4.3), Inches(0.4), trad, size=11, color=C_TEXT_MUTED, align=PP_ALIGN.CENTER)
        
        # Kloudera
        add_card(slide, Inches(8.0), curr_top, Inches(4.5), Inches(0.8), bg_color=(239, 246, 255))
        add_textbox(slide, Inches(8.1), curr_top + Inches(0.2), Inches(4.3), Inches(0.4), kdr, size=11, bold=True, color=C_ELEC_BLUE, align=PP_ALIGN.CENTER)
        
    slide.notes_slide.notes_text_frame.text = (
        "We are positioned as a superior alternative to traditional IT outsourcing. "
        "By integrating security, AI customization, and automation speed, we offer higher value at lower complexity."
    )

    # -------------------------------------------------------------------------
    # SLIDE 13: Founders & Leadership
    # -------------------------------------------------------------------------
    slide = add_content_slide(prs, "Founders & Leadership")
    
    # Left Founder: Ratnesh Pandey
    add_card(slide, Inches(0.8), Inches(1.6), Inches(5.6), Inches(4.8), bg_color=(248, 250, 252))
    if os.path.exists(CEO1_PATH):
        slide.shapes.add_picture(CEO1_PATH, Inches(1.1), Inches(1.9), width=Inches(1.8), height=Inches(1.8))
    else:
        # Placeholder circle
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.1), Inches(1.9), Inches(1.8), Inches(1.8))
        circ.fill.solid()
        circ.fill.fore_color.rgb = RGBColor(*C_DARK_BLUE)
        add_textbox(slide, Inches(1.1), Inches(2.6), Inches(1.8), Inches(0.4), "RP", size=24, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        
    add_textbox(slide, Inches(3.1), Inches(1.9), Inches(3.0), Inches(0.4), "RATNESH PANDEY", size=16, bold=True, color=C_TEXT_DARK)
    add_textbox(slide, Inches(3.1), Inches(2.3), Inches(3.0), Inches(0.4), "CEO | Co-Founder | vCISO", size=12, color=C_CYAN, italic=True)
    add_textbox(slide, Inches(1.1), Inches(3.9), Inches(5.0), Inches(2.2),
                "• Virtual CISO and technology architect for enterprise clients.\n"
                "• Over 12+ years of cybersecurity and infrastructure design experience.\n"
                "• Built custom security systems for global financial technology and IT corporations.",
                size=11, color=C_TEXT_DARK)

    # Right Founder: Deepa Pandey
    add_card(slide, Inches(6.8), Inches(1.6), Inches(5.6), Inches(4.8), bg_color=(248, 250, 252))
    if os.path.exists(CEO2_PATH):
        slide.shapes.add_picture(CEO2_PATH, Inches(7.1), Inches(1.9), width=Inches(1.8), height=Inches(1.8))
    else:
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.1), Inches(1.9), Inches(1.8), Inches(1.8))
        circ.fill.solid()
        circ.fill.fore_color.rgb = RGBColor(*C_DARK_BLUE)
        add_textbox(slide, Inches(7.1), Inches(2.6), Inches(1.8), Inches(0.4), "DP", size=24, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        
    add_textbox(slide, Inches(9.1), Inches(1.9), Inches(3.0), Inches(0.4), "DEEPA PANDEY", size=16, bold=True, color=C_TEXT_DARK)
    add_textbox(slide, Inches(9.1), Inches(2.3), Inches(3.0), Inches(0.4), "Director | Co-Founder | CHRO", size=12, color=C_CYAN, italic=True)
    add_textbox(slide, Inches(7.1), Inches(3.9), Inches(5.0), Inches(2.2),
                "• Strategic management, operations, and human resources leader.\n"
                "• Architect of Kloudera's organizational development and business growth strategy.\n"
                "• Driving local talent engagement and technology enablement initiatives.",
                size=11, color=C_TEXT_DARK)
                
    slide.notes_slide.notes_text_frame.text = (
        "Our leadership team combines deep technical expertise with executive management skills. "
        "CEO Ratnesh Pandey has extensive experience as a Virtual CISO, and Director Deepa Pandey leads HR, "
        "operations, and strategic planning."
    )

    # -------------------------------------------------------------------------
    # SLIDE 14: Strategic Partnerships
    # -------------------------------------------------------------------------
    slide = add_content_slide(prs, "Strategic Alliances")
    
    # 4 partners
    partners = [
        ("MICROSOFT", "Ecosystem Integration", "Full integration with Microsoft Azure, Microsoft Entra Identity governance, and productivity infrastructures."),
        ("MONGODB", "Database Scaling", "Strategic database development and scaling partner, optimizing data models and application queries."),
        ("CROSS CIPHER", "Enterprise Threat Intel", "Providing deep network intelligence, secure corporate gateway solutions, and cloud defenses."),
        ("SPRINTO", "Compliance Automation", "Enabling rapid compliance audits (SOC 2, ISO 27001, HIPAA) and monitoring posture automation.")
    ]
    
    col_w = Inches(2.7)
    for i, (p_name, p_cat, p_desc) in enumerate(partners):
        left = Inches(0.8) + i * (col_w + Inches(0.3))
        top = Inches(1.8)
        card_h = Inches(4.5)
        
        add_card(slide, left, top, col_w, card_h, bg_color=C_WHITE)
        
        # Header box representing logo placement
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, col_w, Inches(1.2))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(*C_DARK_BLUE)
        box.line.fill.background()
        
        add_textbox(slide, left + Inches(0.1), top + Inches(0.4), col_w - Inches(0.2), Inches(0.4),
                    p_name, size=14, bold=True, color=C_CYAN, align=PP_ALIGN.CENTER)
        
        add_textbox(slide, left + Inches(0.2), top + Inches(1.5), col_w - Inches(0.4), Inches(0.4),
                    p_cat, size=11, bold=True, color=C_TEXT_DARK, align=PP_ALIGN.CENTER)
        add_textbox(slide, left + Inches(0.2), top + Inches(2.0), col_w - Inches(0.4), Inches(2.2),
                    p_desc, size=11, color=C_TEXT_MUTED)
                    
    slide.notes_slide.notes_text_frame.text = (
        "We collaborate with global leaders. Our partnerships with Microsoft, MongoDB, Cross Cipher, "
        "and Sprinto allow us to deliver fully integrated, verified, and compliant infrastructures."
    )

    # -------------------------------------------------------------------------
    # SLIDE 15: Client Portfolio
    # -------------------------------------------------------------------------
    slide = add_content_slide(prs, "Trusted Client Portfolio")
    
    # 5 clients
    clients = [
        ("INDIACAPITAL", "FinTech", "Custom enterprise security architectures, financial transaction protection, and threat monitoring."),
        ("FIN CHIKITSAK", "FinTech & Wellness", "Comprehensive software development and advisory portal for financial wellness and advisory."),
        ("GLEEDS", "Global Infrastructure", "International property and construction management consultancy cloud infrastructure scaling."),
        ("STATUSNEO", "Enterprise IT Consulting", "MNC IT consultancy custom web, mobile, and microservices database optimization."),
        ("SIT PUNE", "Education & Research", "Symbiosis Institute of Technology partnership for technical labs, research, and talent acquisition.")
    ]
    
    col_w = Inches(2.1)
    for i, (cl_name, cl_sec, cl_desc) in enumerate(clients):
        left = Inches(0.8) + i * (col_w + Inches(0.2))
        top = Inches(1.8)
        card_h = Inches(4.5)
        
        add_card(slide, left, top, col_w, card_h, bg_color=(250, 252, 255))
        
        # Header banner
        hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, col_w, Inches(0.9))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = RGBColor(*C_DARK_BLUE)
        hdr.line.fill.background()
        
        add_textbox(slide, left + Inches(0.05), top + Inches(0.2), col_w - Inches(0.1), Inches(0.3),
                    cl_name, size=10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, left + Inches(0.05), top + Inches(0.55), col_w - Inches(0.1), Inches(0.3),
                    cl_sec, size=8, color=C_CYAN, align=PP_ALIGN.CENTER, italic=True)
        
        add_textbox(slide, left + Inches(0.15), top + Inches(1.1), col_w - Inches(0.3), Inches(3.2),
                    cl_desc, size=11, color=C_TEXT_DARK)
                    
    slide.notes_slide.notes_text_frame.text = (
        "We are a trusted engineering partner. Our client portfolio spans FinTech, construction consultancy, "
        "enterprise IT, and higher education. This validates the versatility of our solutions layer."
    )

    # -------------------------------------------------------------------------
    # SLIDE 16: Competitive Advantages (Proprietary IP)
    # -------------------------------------------------------------------------
    slide = add_content_slide(prs, "Competitive Advantage: Proprietary IP")
    
    # 3 in-house products
    products = [
        ("KLOUDERA MEET SCHEDULER", "Workspace & Collaboration", 
         "Automated meeting scheduling, multi-timezone synchronization, and encrypted virtual room access governance for enterprise teams."),
        ("KLOUDERA DATA RECOVERY TOOL", "Disaster Recovery & Backup", 
         "Enterprise snapshot restoration, zero-data-loss ransomware recovery, and sector-level disk imaging for cloud and local infrastructure."),
        ("KLOUDERA REMOTE DEVICE CONTROLLER", "Device Management & MDM", 
         "Centralized remote device administration, automated endpoint patch deployment, and real-time security telemetry monitoring.")
    ]
    
    col_w = Inches(3.6)
    for i, (prod_name, prod_cat, prod_desc) in enumerate(products):
        left = Inches(0.8) + i * (col_w + Inches(0.4))
        top = Inches(1.8)
        card_h = Inches(4.5)
        
        add_card(slide, left, top, col_w, card_h, bg_color=C_WHITE)
        
        # Indicator badge
        badge_text = "BETA TESTING" if i==2 else ("IN STAGING" if i==1 else "UNDER DEVELOPMENT")
        
        add_textbox(slide, left + Inches(0.2), top + Inches(0.3), col_w - Inches(0.4), Inches(0.4),
                    prod_name, size=12, bold=True, color=C_DARK_BLUE)
        add_textbox(slide, left + Inches(0.2), top + Inches(0.7), col_w - Inches(0.4), Inches(0.3),
                    prod_cat.upper(), size=9, bold=True, color=C_ELEC_BLUE)
                    
        # Rounded badge box
        badge_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.2), top + Inches(1.1), Inches(1.6), Inches(0.3))
        badge_box.fill.solid()
        badge_box.fill.fore_color.rgb = RGBColor(*C_CYAN)
        badge_box.line.fill.background()
        
        add_textbox(slide, left + Inches(0.2), top + Inches(1.12), Inches(1.6), Inches(0.25),
                    badge_text, size=8, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        
        add_textbox(slide, left + Inches(0.2), top + Inches(1.7), col_w - Inches(0.4), Inches(2.5),
                    prod_desc, size=11, color=C_TEXT_MUTED)
                    
    slide.notes_slide.notes_text_frame.text = (
        "Unlike standard IT consulting firms, Kloudera is building proprietary software products. "
        "Our data recovery, meet scheduler, and remote device manager products are currently in active staging/beta phases, providing high-margin software SaaS potentials."
    )

    # -------------------------------------------------------------------------
    # SLIDE 17: Future Roadmap
    # -------------------------------------------------------------------------
    slide = add_content_slide(prs, "Strategic Roadmap (2026 - 2030)")
    
    # Timeline points
    roadmap = [
        ("2026", "Proprietary IP Rollout", "Release final versions of the Meet Scheduler, Data Recovery Tool, and Remote Controller."),
        ("2027", "Enterprise SaaS Suite", "Launch subscription-based licensing models for our software tools globally."),
        ("2028", "Global Expansion", "Establish operations and sales centers in the Middle East and North America."),
        ("2029", "Strategic Alliances", "Form alliances with tier-1 global system integrators to scale distribution."),
        ("2030", "Vision 2030 Targets", "Establish market leadership in enterprise AI integration and integrated cyber resilience.")
    ]
    
    col_w = Inches(2.2)
    for i, (yr, title_rm, desc_rm) in enumerate(roadmap):
        left = Inches(0.8) + i * (col_w + Inches(0.15))
        top = Inches(2.5)
        card_h = Inches(3.2)
        
        add_card(slide, left, top, col_w, card_h, bg_color=(245, 247, 251))
        
        add_textbox(slide, left + Inches(0.1), top + Inches(0.3), col_w - Inches(0.2), Inches(0.4),
                    yr, size=16, bold=True, color=C_ELEC_BLUE, align=PP_ALIGN.CENTER)
        add_textbox(slide, left + Inches(0.1), top + Inches(0.8), col_w - Inches(0.2), Inches(0.4),
                    title_rm, size=10, bold=True, color=C_TEXT_DARK, align=PP_ALIGN.CENTER)
        add_textbox(slide, left + Inches(0.1), top + Inches(1.3), col_w - Inches(0.2), Inches(1.8),
                    desc_rm, size=10, color=C_TEXT_MUTED, align=PP_ALIGN.CENTER)
                    
    # Connecting line behind the cards
    line_rm = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.0), Inches(11.5), Inches(0.04))
    line_rm.fill.solid()
    line_rm.fill.fore_color.rgb = RGBColor(*C_CYAN)
    line_rm.line.fill.background()
    
    slide.notes_slide.notes_text_frame.text = (
        "Our roadmap is designed for global scale. In 2026, we complete our proprietary tools. "
        "By 2027, we launch SaaS licensing, followed by geographical expansion in 2028."
    )

    # -------------------------------------------------------------------------
    # SLIDE 18: Economic & Social Impact
    # -------------------------------------------------------------------------
    slide = add_content_slide(prs, "Economic & Social Impact")
    
    # 4 impact pillars
    impacts = [
        ("EMPLOYMENT GENERATION", "Expanding engineering hubs in Pune and Bengaluru, providing high-skilled jobs for regional tech talent."),
        ("DIGITAL INDIA SUPPORT", "Aligned with national initiatives, enabling small and mid-sized enterprises (MSMEs) to transition safely to the cloud."),
        ("AI ADOPTION ENABLER", "Sponsoring R&D and training at academic institutes (SIT Pune) to build future-ready AI engineers."),
        ("CYBER SAFETY ADVOCACY", "Conducting security awareness workshops for local communities and startups to prevent digital fraud.")
    ]
    
    col_w = Inches(2.7)
    for i, (im_title, im_desc) in enumerate(impacts):
        left = Inches(0.8) + i * (col_w + Inches(0.3))
        top = Inches(1.8)
        card_h = Inches(4.5)
        
        add_card(slide, left, top, col_w, card_h, bg_color=C_WHITE)
        
        # Icon placeholder box
        ibox = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, col_w, Inches(1.2))
        ibox.fill.solid()
        ibox.fill.fore_color.rgb = RGBColor(*C_DARK_BLUE)
        ibox.line.fill.background()
        
        add_textbox(slide, left + Inches(0.1), top + Inches(0.4), col_w - Inches(0.2), Inches(0.4),
                    "IMPACT", size=11, bold=True, color=C_CYAN, align=PP_ALIGN.CENTER)
        
        add_textbox(slide, left + Inches(0.2), top + Inches(1.5), col_w - Inches(0.4), Inches(0.4),
                    im_title, size=11, bold=True, color=C_TEXT_DARK, align=PP_ALIGN.CENTER)
        add_textbox(slide, left + Inches(0.25), top + Inches(2.1), col_w - Inches(0.5), Inches(2.1),
                    im_desc, size=11, color=C_TEXT_MUTED)
                    
    slide.notes_slide.notes_text_frame.text = (
        "Our business model drives substantial economic and social benefits. "
        "We are contributing directly to national skilling, employment, and digital safety initiatives in India."
    )

    # -------------------------------------------------------------------------
    # SLIDE 19: Vision 2030
    # -------------------------------------------------------------------------
    slide = add_dark_slide(prs, "VISION 2030: GLOBAL RESILIENCE LEADER", 
                           "Pillars for the Next Decade of Enterprise Innovation")
    
    # Three horizontal columns of targets
    targets = [
        ("GLOBAL SCALING", "Establishing sales and distribution in UAE, Singapore, UK, and US markets."),
        ("SaaS DOMINANCE", "Building recurring subscription revenues via proprietary workspace and data tools."),
        ("CYBER RESILIENCE", "Setting international standards for integrated security-by-design structures.")
    ]
    
    col_w = Inches(3.6)
    for i, (tg_title, tg_desc) in enumerate(targets):
        left = Inches(1.0) + i * (col_w + Inches(0.4))
        top = Inches(4.5)
        card_h = Inches(2.0)
        
        add_card(slide, left, top, col_w, card_h, bg_color=C_DARK_BLUE, border_color=C_CYAN)
        
        add_textbox(slide, left + Inches(0.2), top + Inches(0.2), col_w - Inches(0.4), Inches(0.4),
                    tg_title, size=13, bold=True, color=C_WHITE)
        add_textbox(slide, left + Inches(0.2), top + Inches(0.7), col_w - Inches(0.4), Inches(1.1),
                    tg_desc, size=11, color=C_CYAN)
                    
    slide.notes_slide.notes_text_frame.text = (
        "By 2030, Kloudera expects to be a recognized leader in international enterprise resilience. "
        "We will drive SaaS solutions across global markets while maintaining our zero-trust engineering excellence."
    )

    # -------------------------------------------------------------------------
    # SLIDE 20: Closing Slide
    # -------------------------------------------------------------------------
    slide = add_dark_slide(prs, "KLOUDERA TECHNOLOGIES", "Empowering Secure Digital Transformation")
    
    # Left info block
    add_textbox(slide, Inches(1.0), Inches(4.5), Inches(5.0), Inches(2.2),
                "Kloudera Technologies Pvt. Ltd.\n"
                "Web: www.kloudera.ai\n"
                "Email: info@kloudera.ai\n"
                "Tel: +91 9899822926, +91 9390484180",
                size=12, color=C_TEXT_DARK)
                
    # Right info block
    add_textbox(slide, Inches(7.0), Inches(4.5), Inches(5.0), Inches(2.2),
                "Bengaluru Office: Shri Durga Kurupa, Thanisandra, Bengaluru - 560045\n"
                "Pune Office: City Vista DownTown, Kharadi, Pune - 411014\n"
                "Thank You for your consideration.",
                size=12, color=C_ELEC_BLUE)
                
    slide.notes_slide.notes_text_frame.text = (
        "Thank you, judges. We are open for any questions regarding our technologies, market expansion plans, and academic partnerships."
    )

    # Save Presentation
    prs.save('Kloudera_CII_Pitch_Deck_2026.pptx')
    print("Successfully generated Kloudera_CII_Pitch_Deck_2026.pptx")

if __name__ == '__main__':
    generate_deck()
